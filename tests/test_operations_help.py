import unittest
from test_support import TestSupport
import requests
import logging
from pptx import Presentation
from maskgen.ui.help_tools import *
from maskgen.support import *



class TestOperationsHelp(TestSupport):
    def pull_operations_powerpoint(self):
        """
        Will pull latest powerpoint from s3, write to the file in resources.
        :return:
        """
        downloadLink = "https://s3.amazonaws.com/medifor/browser/journal/JournalingToolOperationsDictionary.pptx"
        powerpointPlace = self.locateFile("resources/operationSlides.pptx")

        r = requests.get(downloadLink)
        with open(powerpointPlace, 'wb+') as fd:
            for chunk in r.iter_content(chunk_size=128):
                fd.write(chunk)
        fd.close()
        return powerpointPlace

    def test_operations_help(self):

        powerpointPlace = self.pull_operations_powerpoint() #get latest powerpoint

        operations_file = self.locateFile("resources/operations.json")

        prs = Presentation(powerpointPlace)
        prs.save('Operations.pptx')
        self.addFileToRemove('Operations.pptx')

        ppt_total_slides = len(prs.slides)
        print 'Number of slides gotten from online: ' + str(ppt_total_slides)

        help_loader = HelpLoader()

        self.assertEqual(len(help_loader.missing), 0) #any files missing that were referenced in the loader

        with open(operations_file) as f2:
            operations = json.load(f2)

        #Do all ops have help sections in the linker, and do all help sections point to valid ops?
        opNames_jt = [getValue(op, 'name') for op in getValue(operations, 'operations')]
        opNames_help = getValue(help_loader.linker, 'operation').keys()

        missing_help = [name for name in opNames_jt if name not in opNames_help]
        missing_ops = [name for name in opNames_help if name not in opNames_jt]

        if len(missing_help) > 0:
            logging.getLogger('maskgen').warning('the following operations are not accounted for in the image_linker: ')
            logging.getLogger('maskgen').warning(missing_help)
            raise ValueError('operations missing help.')

        if len(missing_ops) > 0:
            logging.getLogger('maskgen').warning('the following operations are found in the image_linker '
                                                 'but are not found in the operations dictionary: ')
            logging.getLogger('maskgen').warning(missing_ops)
            raise ValueError('invalid/extra operations in help.')

        self.remove_files()

    def test_semantic_group_slides(self):
        downloadLink = "https://s3.amazonaws.com/medifor/browser/journal/SemanticGroups.pptx"
        powerpointPlace = self.locateFile("resources/semanticGroups.pptx")

        r = requests.get(downloadLink)
        with open(powerpointPlace, 'wb+') as fd:
            for chunk in r.iter_content(chunk_size=128):
                fd.write(chunk)
        fd.close()


        imgLinker = self.locateFile("resources/help/image_linker.json")
        groups = self.locateFile("resources/project_properties.json")

        prs = Presentation(powerpointPlace)
        prs.save('Operations.pptx')
        self.addFileToRemove('Operations.pptx')

        prs = Presentation(powerpointPlace)
        prs.save('Operations2.pptx')
        self.addFileToRemove('Operations2.pptx')

        slide = len(prs.slides)
        print 'Number of slides gotten from online: ' + str(slide)

        jtLocation = os.path.join(os.path.split(imgLinker)[0], 'semanticSlides')
        path, dirs, slides = next(os.walk(jtLocation))
        print "JT Semantic Slides: " + str(len(slides))


        with open(imgLinker) as f:
            data = json.load(f)
        with open(groups) as f2:
            semGroups = json.load(f2)
        semanticGroups = []
        for d in semGroups["properties"]:
            try:
                if d["semanticgroup"]:
                    semanticGroups.append(d)
            except:
                pass

        images = set()
        missing = []
        for d in semanticGroups:
            g = d["description"]
            if g not in data["semanticgroup"] or len(data["semanticgroup"][g]["images"]) ==0 or data["semanticgroup"][g]["images"][0] =="":
                missing.append(g)
            else:
                for i in data["semanticgroup"][g]["images"]:
                    self.assertTrue(os.path.exists(os.path.join(jtLocation, os.path.split(i)[1])),
                                    os.path.join(jtLocation, i))
                    images.add(i)
                data["semanticgroup"].pop(g)

        self.assertTrue(missing == [], "Missing is not empty " + str(missing))
        self.assertTrue(len(data["semanticgroup"]) == 0, "There are extra operation(s) in the help section")

        self.remove_files()
        

if __name__ == '__main__':
    unittest.main()
